# -*- coding: utf-8 -*-
"""
文件系统操作模块 (v5.4.3 - Bug 修复)。

此版本修复了 `get_content_slice` 方法中的一个严重 Bug。
原先的 f-string 中错误地使用了双反斜杠 (\\n)，导致返回的
是字面量 '\n' 而不是一个真正的换行符。此版本已将其修正为单反斜杠 (\n)。
"""

import hashlib
import logging
import os
import re
from typing import Iterator

# --- 引入所有需要的第三方文档解析库 ---
import docx
import fitz  # PyMuPDF，用于解析PDF
import openpyxl
import pptx
import xlrd


def _clean_text(text: str) -> str:
    """
    对文本进行清洗，为分词和向量化做准备。
    """
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r"""[^\u4e00-\u9fa5a-zA-Z0-9,.!?;:()"\'[\]]""", ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def scan_files(root_path: str, allowed_extensions: set[str]) -> Iterator[str]:
    """
    递归扫描指定目录下所有符合扩展名要求的文件。
    """
    if not os.path.isdir(root_path):
        logging.warning(f"指定的扫描路径不是一个有效目录: {root_path}")
        return

    for dirpath, _, filenames in os.walk(root_path):
        for filename in filenames:
            if filename.startswith('~$'):
                continue

            file_ext = os.path.splitext(filename)[1].lower()
            if file_ext in allowed_extensions:
                full_path = os.path.join(dirpath, filename)
                normalized_path = full_path.replace('\\', '/')
                yield normalized_path


def calculate_file_hash(file_path: str) -> str | None:
    """
    计算单个文件的 SHA-256 哈希值。
    """
    norm_path = os.path.normpath(file_path)
    sha256_hash = hashlib.sha256()
    try:
        with open(norm_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096 * 1024), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()
    except (IOError, PermissionError) as e:
        logging.error(f"无法读取文件或计算哈希值: {norm_path}, 错误: {e}")
        return None


def calculate_content_hash(content: str) -> str:
    """
    计算字符串内容的 SHA-256 哈希值。
    """
    sha256_hash = hashlib.sha256()
    sha256_hash.update(content.encode('utf-8'))
    return sha256_hash.hexdigest()


def get_content_slice(file_path: str) -> str:
    """
    提取、清洗并返回一个文档的三段式内容摘要。
    """
    norm_path = os.path.normpath(file_path)
    file_ext = os.path.splitext(norm_path)[1].lower()
    text_content = ""

    try:
        if file_ext in ('.txt', '.md'):
            with open(norm_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
        elif file_ext == '.pdf':
            with fitz.open(norm_path) as doc:
                for page in doc:
                    text_content += page.get_text("text", sort=True)
        elif file_ext == '.docx':
            doc = docx.Document(norm_path)
            for para in doc.paragraphs:
                text_content += para.text + '\n'
        elif file_ext == '.pptx':
            prs = pptx.Presentation(norm_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + '\n'
        elif file_ext == '.xlsx':
            workbook = openpyxl.load_workbook(norm_path, read_only=True)
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text_content += str(cell.value) + ' '
                    text_content += '\n'
        elif file_ext == '.xls':
            workbook = xlrd.open_workbook(norm_path)
            for sheet in workbook.sheets():
                for row_idx in range(sheet.nrows):
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        if cell_value:
                            text_content += str(cell_value) + ' '
                    text_content += '\n'
        elif file_ext == '.ppt':
            logging.warning(
                f"'.ppt' (旧版PowerPoint) 文件是二进制格式，当前版本无法直接提取其文本内容。将跳过文件: {norm_path}")
            return ""
    except Exception as e:
        logging.error(f"无法从文件提取文本内容: {norm_path}, 错误: {e}")
        return ""

    cleaned_text = _clean_text(text_content)
    total_len = len(cleaned_text)
    part_size = 2 * 1024

    if total_len <= 3 * part_size:
        return cleaned_text

    head = cleaned_text[:part_size]

    middle_start = (total_len - part_size) // 2
    middle_end = middle_start + part_size
    middle = cleaned_text[middle_start:middle_end]

    tail = cleaned_text[-part_size:]

    # v5.4.3 Bug 修复: 使用单反斜杠 \n 来表示真正的换行符
    return f"{head}\n... (中间部分) ...\n{middle}\n... (结尾部分) ...\n{tail}"
